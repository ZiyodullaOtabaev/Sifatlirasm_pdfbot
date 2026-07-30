"""
AI Presentation / Slide Generator using python-pptx.
Supports 10 Prestigious Academic Template Themes (Harvard, Oxford, Cambridge, Nature, Stanford),
Deep AI Content Generation (Llama-3/Gemini), FLUX AI Real Image Generation & Insertion,
Assertion-Evidence (A-E) Model Headlining, Academic Header & Footer Ribbons,
Standard 12-Slide Deck architecture, and PPTX to PDF conversion.
"""
import os
import uuid
import json
import logging
from datetime import datetime
from typing import List, Dict, Tuple, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from PIL import Image, ImageDraw, ImageFont

logger = logging.getLogger(__name__)

# 10 Prestigious Academic Template Themes (Background, Primary Title, Accent, Card BG, Body Text, Secondary Accent)
SLIDE_THEMES = {
    "oxford_navy": {
        "name": "🏛 Oxford Navy & Gold",
        "bg": RGBColor(15, 34, 64),        # Deep Oxford Navy
        "title": RGBColor(255, 255, 255),  # Pure White
        "accent": RGBColor(212, 175, 55),   # Academic Metallic Gold
        "card_bg": RGBColor(27, 54, 93),   # Oxford Slate Blue Card
        "body": RGBColor(226, 232, 240),   # Light Slate Text
        "badge_bg": RGBColor(40, 75, 120),
    },
    "harvard_crimson": {
        "name": "🎓 Harvard Crimson & Ivory",
        "bg": RGBColor(250, 249, 246),    # Academic Ivory / Off-white
        "title": RGBColor(165, 28, 48),    # Harvard Crimson Maroon
        "accent": RGBColor(30, 58, 138),   # Navy Accent
        "card_bg": RGBColor(255, 255, 255), # Pure White Card
        "body": RGBColor(51, 65, 85),     # Deep Charcoal Body
        "badge_bg": RGBColor(241, 245, 249),
    },
    "cambridge_slate": {
        "name": "🔬 Cambridge Slate & Cyan",
        "bg": RGBColor(15, 23, 42),       # Deep Slate
        "title": RGBColor(248, 250, 252),  # Light Slate
        "accent": RGBColor(56, 189, 248),   # Tech Electric Cyan
        "card_bg": RGBColor(30, 41, 59),   # Dark Card
        "body": RGBColor(203, 213, 225),   # Soft Slate Text
        "badge_bg": RGBColor(51, 65, 85),
    },
    "nature_emerald": {
        "name": "🌿 Nature Emerald & Silver",
        "bg": RGBColor(6, 78, 59),         # Deep Emerald
        "title": RGBColor(254, 240, 138),  # Pale Gold Title
        "accent": RGBColor(234, 179, 8),   # Rich Gold Accent
        "card_bg": RGBColor(4, 120, 87),   # Forest Green Card
        "body": RGBColor(236, 253, 245),   # Soft Mint Body
        "badge_bg": RGBColor(6, 95, 70),
    },
    "stanford_navy": {
        "name": "🏢 Stanford Executive",
        "bg": RGBColor(241, 245, 249),    # Crisp Light Gray
        "title": RGBColor(15, 23, 42),     # Dark Navy
        "accent": RGBColor(37, 99, 235),   # Royal Blue Accent
        "card_bg": RGBColor(255, 255, 255), # Pure White Card
        "body": RGBColor(51, 65, 85),     # Dark Slate Text
        "badge_bg": RGBColor(226, 232, 240),
    },
    "mit_quantum": {
        "name": "⚡️ MIT Quantum Blue",
        "bg": RGBColor(49, 46, 129),       # Deep Quantum Indigo
        "title": RGBColor(255, 255, 255),  # White
        "accent": RGBColor(168, 85, 247),  # Quantum Purple
        "card_bg": RGBColor(67, 56, 202),  # Indigo Card
        "body": RGBColor(224, 231, 255),   # Ice Blue Body
        "badge_bg": RGBColor(79, 70, 229),
    },
    "sorbonne_sand": {
        "name": "☕️ Sorbonne Terracotta",
        "bg": RGBColor(254, 243, 199),    # Warm Sand
        "title": RGBColor(120, 53, 15),    # Terracotta Title
        "accent": RGBColor(154, 52, 18),   # Burnt Orange
        "card_bg": RGBColor(255, 251, 235), # Cream White Card
        "body": RGBColor(69, 26, 3),      # Espresso Body
        "badge_bg": RGBColor(253, 230, 138),
    },
    "imperial_gold": {
        "name": "👑 Imperial Obsidian",
        "bg": RGBColor(9, 9, 11),          # Obsidian Black
        "title": RGBColor(234, 179, 8),    # Pure Gold Title
        "accent": RGBColor(250, 204, 21),  # Bright Gold Accent
        "card_bg": RGBColor(24, 24, 27),   # Obsidian Card
        "body": RGBColor(244, 244, 245),   # White Body
        "badge_bg": RGBColor(39, 39, 42),
    },
    "hopkins_teal": {
        "name": "🏥 Johns Hopkins Medical",
        "bg": RGBColor(236, 254, 255),    # Medical Cyan
        "title": RGBColor(15, 118, 110),   # Deep Teal
        "accent": RGBColor(13, 148, 136),  # Bright Teal Accent
        "card_bg": RGBColor(255, 255, 255), # White Card
        "body": RGBColor(19, 78, 74),     # Dark Teal Body
        "badge_bg": RGBColor(207, 250, 254),
    },
    "eth_mono": {
        "name": "📐 ETH Zurich Mono",
        "bg": RGBColor(255, 255, 255),    # High Contrast White
        "title": RGBColor(0, 0, 0),        # Pitch Black
        "accent": RGBColor(38, 38, 38),    # Dark Charcoal
        "card_bg": RGBColor(245, 245, 245), # Light Gray Card
        "body": RGBColor(23, 23, 23),     # Off-black Body
        "badge_bg": RGBColor(229, 229, 229),
    }
}

THEME_KEYS = list(SLIDE_THEMES.keys())


def generate_deep_ai_slides_content(topic: str) -> Optional[dict]:
    """
    Generate deeply researched, 100% topic-tailored Uzbek academic slide content using AI LLM.
    Returns parsed JSON dictionary or None if generation fails.
    """
    token = os.getenv("REPLICATE_API_TOKEN", "")
    if not token:
        return None

    prompt = f"""You are a distinguished university professor and academic researcher. Create a deeply researched, highly specific 12-slide academic presentation structure in UZBEK language (Latin script) for the topic: "{topic}".

Do NOT use generic text or placeholders. Provide real facts, specific numbers/statistics, real academic laws/sources, and precise terminology about "{topic}".

Return ONLY a raw valid JSON object (no markdown backticks, no text before or after):
{{
  "slide2_text": "2-sentence deep academic overview of {topic}.",
  "slide2_points": ["Fact 1 with details", "Fact 2 with details", "Fact 3 with details"],
  "slide2_metric": {{"title": "🎯 Qamrov", "val": "98.4%", "desc": "Academic metric"}},
  
  "slide3_agenda": [
    ["01", "Subtopic 1"],
    ["02", "Subtopic 2"],
    ["03", "Subtopic 3"],
    ["04", "Subtopic 4"],
    ["05", "Subtopic 5"]
  ],

  "slide4_title": "📚 3. Nazariy Asoslar: Assertion headline about {topic}",
  "slide4_text": "Deep theoretical foundation...",
  "slide4_points": ["Theory concept 1", "Theory concept 2", "Theory concept 3"],
  "slide4_metric": {{"title": "📖 Standart", "val": "IEEE/ISO", "desc": "Academic standard"}},

  "slide5_title": "⚠️ 4. Dolzarb Muammolar: Assertion headline about challenges",
  "slide5_text": "Detailed analysis of current obstacles...",
  "slide5_points": ["Real problem 1", "Real problem 2", "Real problem 3"],
  "slide5_metric": {{"title": "🛑 Risk Tahlili", "val": "High Focus", "desc": "Critical risk"}},

  "slide6_title": "📊 5. Atrofli Tahlil: Assertion headline about stats",
  "slide6_text": "Data analysis breakdown...",
  "slide6_points": ["Specific statistic 1", "Specific statistic 2", "Specific statistic 3"],
  "slide6_metric": {{"title": "📈 O'sish Sur'ati", "val": "+42.5%", "desc": "Metric trend"}},

  "slide7_title": "💼 6. Amaliy Tajriba: Assertion headline about case studies",
  "slide7_text": "Real case studies in Uzbekistan or global practice...",
  "slide7_points": ["Case study 1", "Case study 2", "Case study 3"],
  "slide7_metric": {{"title": "🏢 Amaliyot", "val": "96.8%", "desc": "Implementation rate"}},

  "slide8_title": "⚡️ 7. Asosiy Afzalliklar: Assertion headline about benefits",
  "slide8_text": "Concrete benefits and technological impact...",
  "slide8_points": ["Advantage 1", "Advantage 2", "Advantage 3"],
  "slide8_metric": {{"title": "🚀 Samara", "val": "4.5x Tezroq", "desc": "Efficiency multiplier"}},

  "slide9_title": "💡 8. Strategik Tavsiyalar: Assertion headline about solutions",
  "slide9_text": "Step-by-step action plan and recommendations...",
  "slide9_points": ["Recommendation 1", "Recommendation 2", "Recommendation 3"],
  "slide9_metric": {{"title": "📝 Dastur", "val": "Action Plan", "desc": "Roadmap"}},

  "slide10_title": "🌐 9. Kelajak Istiqbollari: Assertion headline about future",
  "slide10_text": "Future outlook, trends, and AI integration...",
  "slide10_points": ["Future trend 1", "Future trend 2", "Future trend 3"],
  "slide10_metric": {{"title": "🔮 Istiqbol", "val": "Vision 2030", "desc": "Long term target"}},

  "slide11_sources": [
    "1. Specific law or academic journal 1",
    "2. Specific academic book or monograph 2",
    "3. Specific research report or statistical bulletin 3",
    "4. Specific official government or international report 4"
  ]
}}
"""

    try:
        import replicate
        logger.info(f"Calling LLM for deep topic-tailored content: '{topic}'")
        out = replicate.run(
            "meta/meta-llama-3-70b-instruct",
            input={"prompt": prompt, "max_tokens": 2500, "temperature": 0.3}
        )
        raw_text = "".join(out).strip()
        if "```" in raw_text:
            parts = raw_text.split("```")
            raw_text = parts[1] if len(parts) > 1 else raw_text
            if raw_text.startswith("json"):
                raw_text = raw_text[4:].strip()
        
        # Clean control characters that break JSON parsing
        raw_text = raw_text.replace("\n", " ").replace("\r", " ")
        parsed = json.loads(raw_text, strict=False)
        logger.info(f"Successfully generated deep AI content for topic: '{topic}'")
        return parsed
    except Exception as e:
        logger.warning(f"AI content generation fallback: {e}")
        return None


def generate_ai_slide_image(prompt_text: str, download_dir: str = "downloads") -> Optional[str]:
    """
    Generate an HD topic illustration using Replicate API (FLUX/SDXL) or high-res PIL visual artwork card.
    Returns filepath of generated PNG image suitable for python-pptx.
    """
    os.makedirs(download_dir, exist_ok=True)
    out_file = os.path.join(download_dir, f"ai_img_{uuid.uuid4().hex[:8]}.png")

    replicate_token = os.getenv("REPLICATE_API_TOKEN", "")
    if replicate_token:
        try:
            import replicate
            import io
            logger.info(f"Generating FLUX AI image for slide topic: '{prompt_text}'")
            output = replicate.run(
                "black-forest-labs/flux-1-schnell",
                input={"prompt": f"Academic illustration of {prompt_text}, clean vector graphic, professional high quality"}
            )
            if output and len(output) > 0:
                import httpx
                img_url = str(output[0])
                resp = httpx.get(img_url, timeout=20)
                if resp.status_code == 200:
                    img_bytes = io.BytesIO(resp.content)
                    im = Image.open(img_bytes).convert("RGB")
                    im.save(out_file, "PNG")
                    return os.path.abspath(out_file)
        except Exception as e:
            logger.warning(f"Replicate AI image generation fallback: {e}")

    # High-Res Artwork Graphic Card Fallback (Saved as PNG)
    try:
        w, h = 600, 800
        img = Image.new("RGB", (w, h), (30, 41, 59))
        draw = ImageDraw.Draw(img)

        # Draw artistic academic geometry
        draw.rounded_rectangle([30, 30, 570, 770], radius=20, fill=(15, 23, 42), outline=(56, 189, 248), width=4)
        draw.ellipse([150, 150, 450, 450], outline=(168, 85, 247), width=6)
        draw.ellipse([200, 200, 400, 400], outline=(234, 179, 8), width=4)
        
        draw.text((120, 520), "VISUAL TAHLIL", fill=(248, 250, 252))
        draw.text((120, 570), f"Mavzu: {prompt_text[:25]}", fill=(56, 189, 248))
        draw.text((120, 620), "AI Grafik & Metrika", fill=(203, 213, 225))

        img.save(out_file, "PNG")
        return os.path.abspath(out_file)
    except Exception as err:
        logger.error(f"Error creating slide image fallback: {err}")
        return None


def generate_template_preview_image(theme_name: str, download_dir: str = "downloads") -> str:
    """
    Generate a high-res 800x450 visual preview image card of the specified slide template.
    Returns absolute filepath of generated PNG.
    """
    os.makedirs(download_dir, exist_ok=True)
    theme = SLIDE_THEMES.get(theme_name, SLIDE_THEMES["oxford_navy"])

    bg_rgb = tuple(theme["bg"])
    card_rgb = tuple(theme["card_bg"])
    accent_rgb = tuple(theme["accent"])
    title_rgb = tuple(theme["title"])
    body_rgb = tuple(theme["body"])

    width, height = 800, 450
    img = Image.new("RGB", (width, height), bg_rgb)
    draw = ImageDraw.Draw(img)

    # Draw left accent bar
    draw.rectangle([0, 0, 24, height], fill=accent_rgb)

    # Draw main preview card
    draw.rounded_rectangle([60, 45, 740, 405], radius=15, fill=card_rgb, outline=accent_rgb, width=3)

    # Header text & details inside preview card
    draw.text((90, 70), f"Akademik Shablon: {theme['name']}", fill=accent_rgb)
    draw.text((90, 115), "Assertion-Evidence (A-E) Garvard & MIT Standarti", fill=title_rgb)

    # Draw sample bullet points inside preview card
    draw.text((90, 175), "• 12 Betli Standart Akademik & Konferensiya Slaydi", fill=body_rgb)
    draw.text((90, 215), "• Deep Gemini AI Bilan Yozilgan Aniq Faktlar", fill=body_rgb)
    draw.text((90, 255), "• Metrika Kartalari va FLUX AI Illyustratsiyalari", fill=body_rgb)
    draw.text((90, 295), "• Rasmiy Adabiyotlar hamda Manbalar Ro'yxati", fill=body_rgb)

    # Draw right badge box
    draw.rounded_rectangle([520, 165, 710, 335], radius=10, fill=bg_rgb, outline=accent_rgb, width=2)
    draw.text((540, 195), "AKADEMIK", fill=accent_rgb)
    draw.text((540, 235), "12 BET", fill=title_rgb)
    draw.text((540, 275), "IEEE / APA", fill=body_rgb)

    # Footer Ribbon inside preview
    draw.rectangle([60, 370, 740, 405], fill=accent_rgb)
    draw.text((680, 380), "04 / 12", fill=bg_rgb)

    out_filepath = os.path.join(download_dir, f"preview_{theme_name}.png")
    img.save(out_filepath)
    return os.path.abspath(out_filepath)


def create_presentation_slides(
    topic: str,
    theme_name: str = "oxford_navy",
    author_name: str = "",
    institution: str = "",
    download_dir: str = "downloads"
) -> str:
    """
    Generate an ultra-professional 12-slide standard academic presentation PPTX file.
    Follows Harvard & MIT Assertion-Evidence (A-E) model with deep Gemini AI content & FLUX AI images.
    Returns absolute filepath of generated .pptx file.
    """
    os.makedirs(download_dir, exist_ok=True)
    theme = SLIDE_THEMES.get(theme_name, SLIDE_THEMES["oxford_navy"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 Widescreen aspect ratio
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    current_year = datetime.now().year

    author_display = author_name if author_name else "Abdulla Abdullayev"
    institution_display = institution if institution else "Oliy va O'rta Maxsus Ta'lim Muassasasi"

    # Generate deep AI topic-specific presentation content
    ai_content = generate_deep_ai_slides_content(topic)

    # Generate FLUX AI topic illustrations for Slide 4 and Slide 7
    ai_img_1 = generate_ai_slide_image(topic)
    ai_img_2 = generate_ai_slide_image(f"{topic} analysis")

    # Construct 12-Slide Deck with Deep AI Content if available
    if ai_content:
        slide2_data = {
            "text": ai_content.get("slide2_text", f"Ushbu taqdimot '{topic}' mavzusini har tomonlama tahlil qiladi."),
            "points": ai_content.get("slide2_points", ["Mavzuning dolzarbligi", "Asosiy maqsadi", "Kutilayotgan samarasi"]),
            "mtitle": ai_content.get("slide2_metric", {}).get("title", "🎯 Qamrov"),
            "mval": ai_content.get("slide2_metric", {}).get("val", "100%"),
            "mdesc": ai_content.get("slide2_metric", {}).get("desc", "Tahlil darajasi")
        }

        agenda_points = ai_content.get("slide3_agenda", [
            ("01", "Kirish va Nazariy Konsepsiyalar"),
            ("02", "Dolzarb Muammolar va Omil Tahlili"),
            ("03", "Empirik Statistika va Case Study"),
            ("04", "Strategik Tavsiyalar va Yechimlar"),
            ("05", "Akademik Adabiyotlar va Manbalar")
        ])

        slide4_data = {
            "title": ai_content.get("slide4_title", f"📚 3. Nazariy Asoslar: {topic} Konseptsiyalari"),
            "text": ai_content.get("slide4_text", f"'{topic}' nazariy bazasi va ilmiy prinsiplari:"),
            "points": ai_content.get("slide4_points", ["Atamalar ta'rifi", "ISO standartlar", "Ilmiy nazariya"]),
            "mtitle": ai_content.get("slide4_metric", {}).get("title", "📖 Nazariya"),
            "mval": ai_content.get("slide4_metric", {}).get("val", "ISO/IEEE"),
            "mdesc": ai_content.get("slide4_metric", {}).get("desc", "Standartlarga moslik")
        }

        slide5_data = {
            "title": ai_content.get("slide5_title", "⚠️ 4. Dolzarb Muammolar: Mavjud Cheklovlar"),
            "text": ai_content.get("slide5_text", f"Sohadagi asosiy to'siqlar va muammolar:"),
            "points": ai_content.get("slide5_points", ["Resurs cheklovlari", "Infratuzilma", "Risk tahlili"]),
            "mtitle": ai_content.get("slide5_metric", {}).get("title", "🛑 Risk Tahlili"),
            "mval": ai_content.get("slide5_metric", {}).get("val", "High Focus"),
            "mdesc": ai_content.get("slide5_metric", {}).get("desc", "Kritik muammolar")
        }

        slide6_data = {
            "title": ai_content.get("slide6_title", "📊 5. Atrofli Tahlil: Statistik Dinamika"),
            "text": ai_content.get("slide6_text", f"Statistik natijalar va dinamik ko'rsatkichlar:"),
            "points": ai_content.get("slide6_points", ["Dinamik o'sish", "Xarajatlar qisqarishi", "Xalqaro ulush"]),
            "mtitle": ai_content.get("slide6_metric", {}).get("title", "📈 O'sish Sur'ati"),
            "mval": ai_content.get("slide6_metric", {}).get("val", "+48.6%"),
            "mdesc": ai_content.get("slide6_metric", {}).get("desc", "Yillik ko'rsatkich")
        }

        slide7_data = {
            "title": ai_content.get("slide7_title", "💼 6. Amaliy Tajriba: Case Study Tahlili"),
            "text": ai_content.get("slide7_text", f"Amaliyotda joriy etilgan samarali keyslar:"),
            "points": ai_content.get("slide7_points", ["Xalqaro tajriba", "Avtomatlashtirish", "Iqtisodiy samara"]),
            "mtitle": ai_content.get("slide7_metric", {}).get("title", "🏢 Amaliyot"),
            "mval": ai_content.get("slide7_metric", {}).get("val", "99.1%"),
            "mdesc": ai_content.get("slide7_metric", {}).get("desc", "Muvaffaqiyatli joriy etish")
        }

        slide8_data = {
            "title": ai_content.get("slide8_title", "⚡️ 7. Asosiy Afzalliklar: Texnologik Samara"),
            "text": ai_content.get("slide8_text", f"Tizimni tatbiq etishning asosiy afzalliklari:"),
            "points": ai_content.get("slide8_points", ["Tezkorlik", "Resurs tejami", "Ishonchlilik"]),
            "mtitle": ai_content.get("slide8_metric", {}).get("title", "🚀 Samara"),
            "mval": ai_content.get("slide8_metric", {}).get("val", "5x Tezroq"),
            "mdesc": ai_content.get("slide8_metric", {}).get("desc", "Samadorlik ko'rsatkichi")
        }

        slide9_data = {
            "title": ai_content.get("slide9_title", "💡 8. Strategik Tavsiyalar: Harakatlar Dasturi"),
            "text": ai_content.get("slide9_text", f"Sohani rivojlantirish bo'yicha takliflar:"),
            "points": ai_content.get("slide9_points", ["Infratuzilma", "Kadrlar malakasi", "Yo'l xaritasi"]),
            "mtitle": ai_content.get("slide9_metric", {}).get("title", "📝 Dastur"),
            "mval": ai_content.get("slide9_metric", {}).get("val", "Action Roadmap"),
            "mdesc": ai_content.get("slide9_metric", {}).get("desc", "Harakatlar rejasi")
        }

        slide10_data = {
            "title": ai_content.get("slide10_title", "🌐 9. Kelajak Istiqbollari: Global Tendensiyalar"),
            "text": ai_content.get("slide10_text", f"Kelgusi yillardagi rivojlanish vektori:"),
            "points": ai_content.get("slide10_points", ["Raqamli transformatsiya", "Xalqaro hamkorlik", "AI integratsiyasi"]),
            "mtitle": ai_content.get("slide10_metric", {}).get("title", "🔮 Istiqbol"),
            "mval": ai_content.get("slide10_metric", {}).get("val", "Vision 2030"),
            "mdesc": ai_content.get("slide10_metric", {}).get("desc", "Uzoq muddatli maqsad")
        }

        sources_list = ai_content.get("slide11_sources", [
            f"1. O'zbekiston Respublikasi Me me me'yoriy Hujjatlari Portali ({current_year}).",
            f"2. '{topic}' Bo'yicha Xalqaro Ilmiy Jurnallar va Nature/Springer Manbalari.",
            "3. Oliy Ta'lim Muassasalari Akademik Darsliklari hamda Ilmiy To'plamlari.",
            "4. Rasmiy Davlat Statistika Qo'mitasi Hisobotlari."
        ])
    else:
        # Fallback Dynamic Content
        slide2_data = {
            "text": f"Ushbu taqdimot '{topic}' mavzusining fundamental nazariyalari va amaliy samadorligini atrofli tahlil qiladi.",
            "points": [
                f"'{topic}' sohasining strategik dolzarbligi va amaliy ahamiyati.",
                "Tadqiqot va amaliyot yo'nalishidagi ustuvor vazifalar.",
                "Kutilayotgan ijtimoiy-iqtisodiy hamda texnologik samadorlik ko'rsatkichlari."
            ],
            "mtitle": "🎯 Qamrov", "mval": "100%", "mdesc": "Akademik tahlil darajasi"
        }
        agenda_points = [
            ("01", "Kirish va Nazariy Konsepsiyalar"),
            ("02", "Dolzarb Muammolar va Omil Tahlili"),
            ("03", "Empirik Statistika va Case Study"),
            ("04", "Strategik Tavsiyalar va Yechimlar"),
            ("05", "Akademik Adabiyotlar va Manbalar")
        ]
        slide4_data = {
            "title": f"📚 3. Nazariy Asoslar: Fundamental Prinsiplar Barqarorlikni Ta'minlaydi",
            "text": f"'{topic}' tushunchasining nazariy bazasi va fundamental konseptsiyalari:",
            "points": ["Asosiy ilmiy atamalar va ta'riflar", "Xalqaro ISO/IEEE standartlariga moslik", "Nazariyalarning o'zaro integratsiyasi"],
            "mtitle": "📖 Nazariya", "mval": "ISO/IEEE", "mdesc": "Xalqaro me'yorlarga moslik"
        }
        slide5_data = {
            "title": f"⚠️ 4. Dolzarb Muammolar: Resurs va Infratuzilma Cheklovlarini Hal Etish Zarur",
            "text": f"Bugungi kunda '{topic}' sohasida duch kelinayotgan asosiy muammolar:",
            "points": ["Moddiy resurs yetishmovchiligi", "Infratuzilma va kadrlar malakasi", "Risk ko'rsatkichlarini boshqarish"],
            "mtitle": "🛑 Risk Tahlili", "mval": "High Focus", "mdesc": "Kritik muammolar"
        }
        slide6_data = {
            "title": f"📊 5. Atrofli Tahlil: Statistik Ko'rsatkichlar Dinamik O'sishni Ko'rsatmoqda",
            "text": f"So'nggi yillarda '{topic}' yo'nalishi bo'yicha erishilgan statistik natijalar:",
            "points": ["Samadorlikning yillik dinamik o'sishi", "Amaliy jarayonlar tejamkorligi", "Xalqaro ko'rsatkichlar ulushi"],
            "mtitle": "📈 O'sish Sur'ati", "mval": "+48.6%", "mdesc": "Yillik ko'rsatkich"
        }
        slide7_data = {
            "title": f"💼 6. Amaliy Tajriba: Yetakchi Muassasalar Amaliyotida Muvaffaqiyatli Natijalar",
            "text": f"'{topic}' doirasida joriy etilgan eng samarali amaliy keyslar (Case Studies):",
            "points": ["Muvaffaqiyatli xalqaro loyihalar", "Jarayonlarni optimallashtirish", "Amaliy iqtisodiy samara"],
            "mtitle": "🏢 Amaliyot", "mval": "99.1%", "mdesc": "Muvaffaqiyatli keyslar"
        }
        slide8_data = {
            "title": f"⚡️ 7. Asosiy Afzalliklar: Raqamli Texnologiyalar Samadorlikni Oshiradi",
            "text": f"Tizimni tatbiq etish orqali erishiladigan asosiy afzalliklar:",
            "points": ["Tezkor qaror qabul qilish", "Resurs va vaqt tejami", "Tizim barqarorligi va xavfsizlik"],
            "mtitle": "🚀 Samara", "mval": "5x Tezroq", "mdesc": "Samadorlik ko'rsatkichi"
        }
        slide9_data = {
            "title": f"💡 8. Strategik Tavsiyalar: Bosqichma-Bosqich Harakatlar Dasturini Amalga Oshirish",
            "text": f"'{topic}' sohasini rivojlantirish bo'yicha taklif etilayotgan akademik yechimlar:",
            "points": ["Moddiy-texnik bazani zamonaviylashtirish", "Kadrlar malakasini oshirish", "Uzoq muddatli yo'l xaritasi"],
            "mtitle": "📝 Dastur", "mval": "Action Roadmap", "mdesc": "Harakatlar rejasi"
        }
        slide10_data = {
            "title": f"🌐 9. Kelajak Istiqbollari: Global Tendensiyalar va Sun'iy Intellekt Integratsiyasi",
            "text": f"Kelgusi yillarda '{topic}' sohasining rivojlanish vektori:",
            "points": ["Raqamli transformatsiya va AI", "Xalqaro hamkorlik", "Barqaror rivojlanish maqsadlari"],
            "mtitle": "🔮 Istiqbol", "mval": "Vision 2030", "mdesc": "Uzoq muddatli target"
        }
        sources_list = [
            f"1. O'zbekiston Respublikasi Me me me'yoriy Hujjatlari Portali ({current_year}).",
            f"2. '{topic}' Bo'yicha Xalqaro Ilmiy Jurnallar va Nature/Springer Manbalari.",
            "3. Oliy Ta'lim Muassasalari Akademik Darsliklari hamda Ilmiy To'plamlari.",
            "4. Rasmiy Davlat Statistika Qo'mitasi Hisobotlari."
        ]

    slides_data = [
        # Slide 1: Title Cover
        {
            "type": "title_cover",
            "title": topic.title(),
            "author": author_display,
            "institution": institution_display,
            "date": f"{current_year}-yil",
        },
        # Slide 2: Executive Summary
        {
            "type": "content_visual",
            "slide_num": 2,
            "title": f"📑 1. Mavzuning Qisqacha Mazmuni: {topic} Strategik Rivojlanish Negizidir",
            "text": slide2_data["text"],
            "points": slide2_data["points"],
            "metric_title": slide2_data["mtitle"],
            "metric_val": slide2_data["mval"],
            "metric_desc": slide2_data["mdesc"]
        },
        # Slide 3: Agenda & Table of Contents
        {
            "type": "agenda_grid",
            "slide_num": 3,
            "title": "📋 2. Taqdimot Rejasi va Tadqiqot Bosqichlari",
            "points": agenda_points
        },
        # Slide 4: Theoretical Background (Embedded AI Image 1)
        {
            "type": "content_visual",
            "slide_num": 4,
            "title": slide4_data["title"],
            "text": slide4_data["text"],
            "points": slide4_data["points"],
            "image_path": ai_img_1,
            "metric_title": slide4_data["mtitle"],
            "metric_val": slide4_data["mval"],
            "metric_desc": slide4_data["mdesc"]
        },
        # Slide 5: Key Challenges
        {
            "type": "content_visual",
            "slide_num": 5,
            "title": slide5_data["title"],
            "text": slide5_data["text"],
            "points": slide5_data["points"],
            "metric_title": slide5_data["mtitle"],
            "metric_val": slide5_data["mval"],
            "metric_desc": slide5_data["mdesc"]
        },
        # Slide 6: Deep Data Analysis & Metrics
        {
            "type": "content_visual",
            "slide_num": 6,
            "title": slide6_data["title"],
            "text": slide6_data["text"],
            "points": slide6_data["points"],
            "metric_title": slide6_data["mtitle"],
            "metric_val": slide6_data["mval"],
            "metric_desc": slide6_data["mdesc"]
        },
        # Slide 7: Practical Applications & Case Studies (Embedded AI Image 2)
        {
            "type": "content_visual",
            "slide_num": 7,
            "title": slide7_data["title"],
            "text": slide7_data["text"],
            "points": slide7_data["points"],
            "image_path": ai_img_2,
            "metric_title": slide7_data["mtitle"],
            "metric_val": slide7_data["mval"],
            "metric_desc": slide7_data["mdesc"]
        },
        # Slide 8: Key Advantages & Impact
        {
            "type": "content_visual",
            "slide_num": 8,
            "title": slide8_data["title"],
            "text": slide8_data["text"],
            "points": slide8_data["points"],
            "metric_title": slide8_data["mtitle"],
            "metric_val": slide8_data["mval"],
            "metric_desc": slide8_data["mdesc"]
        },
        # Slide 9: Strategic Recommendations & Solutions
        {
            "type": "content_visual",
            "slide_num": 9,
            "title": slide9_data["title"],
            "text": slide9_data["text"],
            "points": slide9_data["points"],
            "metric_title": slide9_data["mtitle"],
            "metric_val": slide9_data["mval"],
            "metric_desc": slide9_data["mdesc"]
        },
        # Slide 10: Future Outlook & Global Trends
        {
            "type": "content_visual",
            "slide_num": 10,
            "title": slide10_data["title"],
            "text": slide10_data["text"],
            "points": slide10_data["points"],
            "metric_title": slide10_data["mtitle"],
            "metric_val": slide10_data["mval"],
            "metric_desc": slide10_data["mdesc"]
        },
        # Slide 11: References & Bibliography
        {
            "type": "references",
            "slide_num": 11,
            "title": "📚 10. Foydalanilgan Akademik Adabiyotlar va Manbalar (IEEE/APA Standard)",
            "sources": sources_list
        },
        # Slide 12: Closing Cover
        {
            "type": "thank_you_clean",
            "slide_num": 12,
            "title": "✨ E'tiboringiz uchun Rahmat!",
            "subtitle": f"Taqdimotchi: {author_display}\nMuassasa: {institution_display}\nMavzu: {topic[:60]}",
        }
    ]

    for s_idx, data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        slide_num = data.get("slide_num", s_idx + 1)
        
        # Background Fill
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = theme["bg"]
        bg_shape.line.fill.background()

        # Add Academic Top Ribbon and Bottom Footer (except title cover)
        if data["type"] != "title_cover":
            # Academic Top Header Line
            top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.04))
            top_line.fill.solid()
            top_line.fill.fore_color.rgb = theme["accent"]
            top_line.line.fill.background()

            # Academic Footer Line
            bot_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.03))
            bot_line.fill.solid()
            bot_line.fill.fore_color.rgb = theme["accent"]
            bot_line.line.fill.background()

            # Academic Footer Text Frame (Only clean page number)
            ft_tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.95), Inches(11.733), Inches(0.4))
            p_ft = ft_tb.text_frame.paragraphs[0]
            p_ft.text = f"{slide_num:02d} / 12"
            p_ft.font.size = Pt(12)
            p_ft.font.bold = True
            p_ft.font.color.rgb = theme["body"]
            p_ft.alignment = PP_ALIGN.RIGHT

        # 1. Title Cover Slide
        if data["type"] == "title_cover":
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), prs.slide_height)
            bar.fill.solid()
            bar.fill.fore_color.rgb = theme["accent"]
            bar.line.fill.background()

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(11.0), Inches(5.1))
            card.fill.solid()
            card.fill.fore_color.rgb = theme["card_bg"]
            card.line.color.rgb = theme["accent"]

            tf = card.text_frame
            tf.word_wrap = True
            
            p_title = tf.paragraphs[0]
            p_title.text = data["title"]
            p_title.font.size = Pt(34)
            p_title.font.bold = True
            p_title.font.color.rgb = theme["accent"]
            p_title.alignment = PP_ALIGN.CENTER
            p_title.space_after = Pt(16)

            p_meta1 = tf.add_paragraph()
            p_meta1.text = f"👤 Taqdimotchi: {data['author']}"
            p_meta1.font.size = Pt(20)
            p_meta1.font.bold = True
            p_meta1.font.color.rgb = theme["title"]
            p_meta1.alignment = PP_ALIGN.CENTER

            p_meta2 = tf.add_paragraph()
            p_meta2.text = f"🏛 Muassasa: {data['institution']}"
            p_meta2.font.size = Pt(18)
            p_meta2.font.color.rgb = theme["body"]
            p_meta2.alignment = PP_ALIGN.CENTER
            p_meta2.space_before = Pt(6)

            p_date = tf.add_paragraph()
            p_date.text = f"📅 Sana: {data['date']}"
            p_date.font.size = Pt(15)
            p_date.font.color.rgb = theme["accent"]
            p_date.alignment = PP_ALIGN.CENTER
            p_date.space_before = Pt(12)

        # 2. Agenda Grid Slide
        elif data["type"] == "agenda_grid":
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.733), Inches(0.8))
            p = tb.text_frame.paragraphs[0]
            p.text = data["title"]
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = theme["accent"]

            top_y = Inches(1.5)
            for idx, (num, item_text) in enumerate(data["points"]):
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_y + Inches(idx * 1.05), Inches(11.733), Inches(0.9))
                card.fill.solid()
                card.fill.fore_color.rgb = theme["card_bg"]
                card.line.color.rgb = theme["accent"]

                tf = card.text_frame
                tf.word_wrap = True
                p_item = tf.paragraphs[0]
                p_item.text = f"{num}.  {item_text}"
                p_item.font.size = Pt(19)
                p_item.font.bold = True
                p_item.font.color.rgb = theme["title"]

        # 3. Content + Visual Illustration Box / AI Image Slide
        elif data["type"] == "content_visual":
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.733), Inches(0.85))
            p = tb.text_frame.paragraphs[0]
            p.text = data["title"]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = theme["accent"]

            # Left Content Card (65% width)
            left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(7.8), Inches(5.1))
            left_card.fill.solid()
            left_card.fill.fore_color.rgb = theme["card_bg"]
            left_card.line.fill.background()

            tf_left = left_card.text_frame
            tf_left.word_wrap = True
            
            p_desc = tf_left.paragraphs[0]
            p_desc.text = data["text"]
            p_desc.font.size = Pt(18)
            p_desc.font.bold = True
            p_desc.font.color.rgb = theme["title"]
            p_desc.space_after = Pt(12)

            for pt in data["points"]:
                p_pt = tf_left.add_paragraph()
                p_pt.text = f"• {pt}"
                p_pt.font.size = Pt(16)
                p_pt.font.color.rgb = theme["body"]
                p_pt.space_before = Pt(8)

            # Right Column: Embed Real AI Image if available, else Metric Card
            img_file = data.get("image_path")
            if img_file and os.path.exists(img_file):
                try:
                    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.5), Inches(3.633), Inches(5.1))
                    frame.fill.solid()
                    frame.fill.fore_color.rgb = theme["card_bg"]
                    frame.line.color.rgb = theme["accent"]

                    slide.shapes.add_picture(img_file, Inches(9.0), Inches(1.6), Inches(3.433), Inches(4.9))
                except Exception as img_err:
                    logger.warning(f"Error adding picture to slide: {img_err}")
            else:
                # Metric Card Fallback
                right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.5), Inches(3.633), Inches(5.1))
                right_card.fill.solid()
                right_card.fill.fore_color.rgb = theme["badge_bg"]
                right_card.line.color.rgb = theme["accent"]

                tf_right = right_card.text_frame
                tf_right.word_wrap = True

                p_mtitle = tf_right.paragraphs[0]
                p_mtitle.text = f"\n{data['metric_title']}"
                p_mtitle.font.size = Pt(22)
                p_mtitle.font.bold = True
                p_mtitle.font.color.rgb = theme["accent"]
                p_mtitle.alignment = PP_ALIGN.CENTER

                p_mval = tf_right.add_paragraph()
                p_mval.text = f"\n{data['metric_val']}"
                p_mval.font.size = Pt(40)
                p_mval.font.bold = True
                p_mval.font.color.rgb = theme["title"]
                p_mval.alignment = PP_ALIGN.CENTER

                p_mdesc = tf_right.add_paragraph()
                p_mdesc.text = f"\n{data['metric_desc']}"
                p_mdesc.font.size = Pt(14)
                p_mdesc.font.color.rgb = theme["body"]
                p_mdesc.alignment = PP_ALIGN.CENTER

        # 4. References & Bibliography Slide
        elif data["type"] == "references":
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.733), Inches(0.8))
            p = tb.text_frame.paragraphs[0]
            p.text = data["title"]
            p.font.size = Pt(25)
            p.font.bold = True
            p.font.color.rgb = theme["accent"]

            ref_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.1))
            ref_card.fill.solid()
            ref_card.fill.fore_color.rgb = theme["card_bg"]
            ref_card.line.fill.background()

            tf_ref = ref_card.text_frame
            tf_ref.word_wrap = True

            for idx, source in enumerate(data["sources"]):
                p_src = tf_ref.paragraphs[0] if idx == 0 else tf_ref.add_paragraph()
                p_src.text = source
                p_src.font.size = Pt(17)
                p_src.font.color.rgb = theme["title"]
                p_src.space_before = Pt(12)

        # 5. Thank You Clean Slide (No bot watermark)
        elif data["type"] == "thank_you_clean":
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(1.8), Inches(9.333), Inches(4.0))
            card.fill.solid()
            card.fill.fore_color.rgb = theme["card_bg"]
            card.line.color.rgb = theme["accent"]

            tf = card.text_frame
            tf.word_wrap = True
            p1 = tf.paragraphs[0]
            p1.text = data["title"]
            p1.font.size = Pt(36)
            p1.font.bold = True
            p1.font.color.rgb = theme["accent"]
            p1.alignment = PP_ALIGN.CENTER

            p2 = tf.add_paragraph()
            p2.text = f"\n{data['subtitle']}"
            p2.font.size = Pt(19)
            p2.font.color.rgb = theme["title"]
            p2.alignment = PP_ALIGN.CENTER

    out_filename = f"slayd_{uuid.uuid4().hex[:8]}.pptx"
    out_filepath = os.path.join(download_dir, out_filename)
    prs.save(out_filepath)
    logger.info(f"Ultra-professional 12-slide presentation generated successfully with theme {theme_name}: {out_filepath}")

    # Cleanup temporary AI images
    for img_p in [ai_img_1, ai_img_2]:
        if img_p and os.path.exists(img_p):
            try:
                os.remove(img_p)
            except Exception:
                pass

    return os.path.abspath(out_filepath)


def convert_pptx_to_pdf(pptx_path: str, theme_name: str = "oxford_navy") -> str:
    """
    Convert presentation PPTX content into a landscape PDF file using ReportLab.
    Returns absolute filepath of generated .pdf file.
    """
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import HexColor

    pdf_path = pptx_path.replace(".pptx", ".pdf")
    theme = SLIDE_THEMES.get(theme_name, SLIDE_THEMES["oxford_navy"])

    # Extract text from PPTX slides
    prs = Presentation(pptx_path)
    
    doc = SimpleDocTemplate(
        pdf_path,
        pagesize=landscape(A4),
        rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36
    )

    styles = getSampleStyleSheet()

    title_hex = f"#{theme['accent'][0]:02x}{theme['accent'][1]:02x}{theme['accent'][2]:02x}"
    text_hex = f"#{theme['title'][0]:02x}{theme['title'][1]:02x}{theme['title'][2]:02x}"

    title_style = ParagraphStyle(
        'SlideTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=24,
        leading=28,
        textColor=HexColor(title_hex),
        spaceAfter=15
    )

    body_style = ParagraphStyle(
        'SlideBody',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=14,
        leading=18,
        textColor=HexColor(text_hex),
        spaceAfter=8
    )

    story = []

    for slide_idx, slide in enumerate(prs.slides):
        story.append(Paragraph(f"<b>Slayd {slide_idx + 1}</b>", title_style))
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        story.append(Paragraph(text, body_style))
        story.append(Spacer(1, 20))

    doc.build(story)
    logger.info(f"Converted PPTX to PDF successfully: {pdf_path}")
    return os.path.abspath(pdf_path)
